"""
Shaman Assistant Document Tools
===============================
Tools for creating and reading documents (PowerPoint, Word, PDF, etc.).
"""
from pathlib import Path
from datetime import datetime
import json
import os

OUTPUT_DIR = Path("./output")
OUTPUT_DIR.mkdir(exist_ok=True)

# Base directory for file reading (project root)
PROJECT_ROOT = Path(__file__).parent.parent.resolve()

# Academy content sync configuration
# Source of truth for curriculum.json and lesson .md files
QUIZ_DEMO_PATH = Path.home() / "PycharmProjects" / "shaman-quiz-demo"
QUIZ_DEMO_CONTENT = QUIZ_DEMO_PATH / "src" / "content"
LOCAL_ACADEMY_PATH = OUTPUT_DIR / "shaman-academy"

# Allowed directories for file reading (security)
ALLOWED_READ_DIRS = [
    PROJECT_ROOT / "output",
    PROJECT_ROOT / "content_input",
    PROJECT_ROOT / "transcripts",
    PROJECT_ROOT / "templates",
    PROJECT_ROOT,  # Allow reading from project root for files like .md, .json, .txt
    QUIZ_DEMO_CONTENT,  # Allow reading from quiz-demo source of truth
]


def create_pptx(title: str, filename: str, slides: list[dict]) -> str:
    """
    Create a PowerPoint presentation.

    Args:
        title: Presentation title
        filename: Output filename (without extension)
        slides: List of slide dicts with:
            - title: str
            - subtitle: str (optional)
            - content: list[str] (bullet points)
            - notes: str (speaker notes, optional)
            - layout: "title" | "section" | "content" | "two_column"

    Returns:
        Path to created file
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TWO_CONTENT = 3

    for i, slide_data in enumerate(slides):
        layout_name = slide_data.get("layout", "content")

        if layout_name == "title" or i == 0:
            layout = prs.slide_layouts[LAYOUT_TITLE]
        elif layout_name == "section":
            layout = prs.slide_layouts[LAYOUT_SECTION]
        elif layout_name == "two_column":
            layout = prs.slide_layouts[LAYOUT_TWO_CONTENT]
        else:
            layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]

        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        if layout_name == "title" and slide_data.get("subtitle"):
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = slide_data["subtitle"]
                    break

        if content := slide_data.get("content"):
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1 and layout_name != "title":
                    tf = shape.text_frame
                    tf.clear()
                    for j, item in enumerate(content):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = item
                        p.level = 0
                    break

        if notes := slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = notes

    filepath = OUTPUT_DIR / f"{filename}.pptx"
    prs.save(str(filepath))
    return str(filepath)


def create_docx(title: str, filename: str, sections: list[dict]) -> str:
    """
    Create a Word document.

    Args:
        title: Document title
        filename: Output filename (without extension)
        sections: List of section dicts with:
            - heading: str
            - level: int (1-3, default 1)
            - content: str (paragraphs, use \\n\\n for breaks)
            - bullets: list[str] (optional)

    Returns:
        Path to created file
    """
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph("")

    for section in sections:
        level = section.get("level", 1)
        doc.add_heading(section["heading"], level=level)

        if content := section.get("content"):
            for para_text in content.split("\n\n"):
                if para_text.strip():
                    doc.add_paragraph(para_text.strip())

        if bullets := section.get("bullets"):
            for bullet in bullets:
                doc.add_paragraph(bullet, style="List Bullet")

    filepath = OUTPUT_DIR / f"{filename}.docx"
    doc.save(str(filepath))
    return str(filepath)


def create_pdf(title: str, filename: str, content: str) -> str:
    """
    Create a PDF from markdown content.

    Args:
        title: Document title
        filename: Output filename (without extension)
        content: Markdown-formatted content

    Returns:
        Path to created file
    """
    try:
        import markdown
        from weasyprint import HTML
    except ImportError:
        # Fallback to markdown file
        filepath = OUTPUT_DIR / f"{filename}.md"
        filepath.write_text(f"# {title}\n\n{content}")
        return str(filepath)

    html_content = markdown.markdown(content, extensions=['tables', 'fenced_code'])

    full_html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>{title}</title>
        <style>
            @page {{ size: A4; margin: 2.5cm; }}
            body {{ font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 11pt; line-height: 1.6; color: #333; }}
            h1 {{ color: #1a365d; border-bottom: 2px solid #1a365d; padding-bottom: 10px; }}
            h2 {{ color: #2c5282; margin-top: 1.5em; }}
            table {{ width: 100%; border-collapse: collapse; margin: 1em 0; }}
            th, td {{ border: 1px solid #e2e8f0; padding: 8px 12px; text-align: left; }}
            th {{ background-color: #edf2f7; }}
            code {{ background-color: #f7fafc; padding: 2px 6px; border-radius: 3px; }}
        </style>
    </head>
    <body>
        <h1>{title}</h1>
        <p style="color: #718096;">Generated: {datetime.now().strftime('%Y-%m-%d')}</p>
        {html_content}
    </body>
    </html>
    """

    filepath = OUTPUT_DIR / f"{filename}.pdf"
    HTML(string=full_html).write_pdf(str(filepath))
    return str(filepath)


def create_markdown(title: str, filename: str, content: str) -> str:
    """
    Create a Markdown file.

    Args:
        title: Document title (used as H1 heading)
        filename: Output filename (without extension)
        content: Markdown-formatted content

    Returns:
        Path to created file
    """
    full_content = f"# {title}\n\n{content}"
    filepath = OUTPUT_DIR / f"{filename}.md"
    filepath.write_text(full_content, encoding="utf-8")
    return str(filepath)


def create_json(filename: str, data: dict | list, schema_description: str = None) -> str:
    """
    Create a JSON file with structured data.

    Args:
        filename: Output filename (without extension)
        data: The JSON data structure (dict or list)
        schema_description: Optional description of the schema used

    Returns:
        Path to created file
    """
    import json

    output = {
        "_meta": {
            "generated": datetime.now().isoformat(),
            "generator": "Shaman Assistant"
        }
    }

    if schema_description:
        output["_meta"]["schema_description"] = schema_description

    # Merge data into output (or wrap if it's a list)
    if isinstance(data, dict):
        output["data"] = data
    else:
        output["data"] = data

    filepath = OUTPUT_DIR / f"{filename}.json"
    filepath.write_text(json.dumps(output, indent=2, ensure_ascii=False), encoding="utf-8")
    return str(filepath)


# =============================================================================
# FILE READING TOOLS
# =============================================================================

def _resolve_file_path(file_path: str) -> Path:
    """
    Resolve a file path, handling both absolute and relative paths.
    For relative paths, resolves against PROJECT_ROOT.

    Args:
        file_path: File path (absolute or relative)

    Returns:
        Resolved absolute Path object

    Raises:
        ValueError: If path is outside allowed directories
    """
    path = Path(file_path)

    # If relative, resolve against project root
    if not path.is_absolute():
        path = PROJECT_ROOT / path

    path = path.resolve()

    # Security check: ensure path is within allowed directories
    is_allowed = any(
        path == allowed_dir or allowed_dir in path.parents or path.is_relative_to(allowed_dir)
        for allowed_dir in ALLOWED_READ_DIRS
    )

    if not is_allowed:
        raise ValueError(f"Access denied: {file_path} is outside allowed directories")

    return path


def read_text_file(file_path: str) -> dict:
    """
    Read a plain text file (.txt, .md, .log, etc.).

    Args:
        file_path: Path to the file (absolute or relative to project root)

    Returns:
        Dict with content and metadata
    """
    try:
        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        content = path.read_text(encoding="utf-8")

        return {
            "file_path": str(path),
            "file_name": path.name,
            "file_type": path.suffix,
            "size_bytes": path.stat().st_size,
            "content": content
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read file: {str(e)}"}


def read_json_file(file_path: str) -> dict:
    """
    Read and parse a JSON file.

    Args:
        file_path: Path to the JSON file

    Returns:
        Dict with parsed data and metadata
    """
    try:
        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        content = path.read_text(encoding="utf-8")
        data = json.loads(content)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "data": data
        }
    except ValueError as e:
        return {"error": str(e)}
    except json.JSONDecodeError as e:
        return {"error": f"Invalid JSON: {str(e)}"}
    except Exception as e:
        return {"error": f"Failed to read file: {str(e)}"}


def read_csv_file(file_path: str, max_rows: int = 100) -> dict:
    """
    Read a CSV file and return as list of dicts.

    Args:
        file_path: Path to the CSV file
        max_rows: Maximum rows to return (default 100)

    Returns:
        Dict with data rows and metadata
    """
    try:
        import csv

        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        with open(path, newline='', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            rows = []
            for i, row in enumerate(reader):
                if i >= max_rows:
                    break
                rows.append(row)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "columns": list(rows[0].keys()) if rows else [],
            "row_count": len(rows),
            "truncated": len(rows) == max_rows,
            "data": rows
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read CSV: {str(e)}"}


def read_excel_file(file_path: str, sheet_name: str = None, max_rows: int = 100) -> dict:
    """
    Read an Excel file (.xlsx, .xls).

    Args:
        file_path: Path to the Excel file
        sheet_name: Specific sheet to read (default: first sheet)
        max_rows: Maximum rows to return (default 100)

    Returns:
        Dict with data and metadata
    """
    try:
        import openpyxl

        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)

        # Get sheet names
        sheet_names = wb.sheetnames

        # Select sheet
        if sheet_name:
            if sheet_name not in sheet_names:
                return {"error": f"Sheet '{sheet_name}' not found. Available: {sheet_names}"}
            ws = wb[sheet_name]
        else:
            ws = wb.active
            sheet_name = ws.title

        # Read data
        rows = []
        headers = None
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i == 0:
                headers = [str(cell) if cell else f"Column_{j}" for j, cell in enumerate(row)]
            else:
                if i > max_rows:
                    break
                row_dict = {headers[j]: cell for j, cell in enumerate(row) if j < len(headers)}
                rows.append(row_dict)

        wb.close()

        return {
            "file_path": str(path),
            "file_name": path.name,
            "sheet_name": sheet_name,
            "all_sheets": sheet_names,
            "columns": headers or [],
            "row_count": len(rows),
            "truncated": len(rows) == max_rows,
            "data": rows
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read Excel: {str(e)}"}


def read_pdf_file(file_path: str, max_pages: int = 20) -> dict:
    """
    Read and extract text from a PDF file.

    Args:
        file_path: Path to the PDF file
        max_pages: Maximum pages to read (default 20)

    Returns:
        Dict with extracted text and metadata
    """
    try:
        import pdfplumber

        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        pages_content = []
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            pages_to_read = min(total_pages, max_pages)

            for i in range(pages_to_read):
                page = pdf.pages[i]
                text = page.extract_text() or ""
                pages_content.append({
                    "page": i + 1,
                    "text": text
                })

        full_text = "\n\n".join(p["text"] for p in pages_content)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "total_pages": total_pages,
            "pages_read": pages_to_read,
            "truncated": pages_to_read < total_pages,
            "content": full_text,
            "pages": pages_content
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read PDF: {str(e)}"}


def read_docx_file(file_path: str) -> dict:
    """
    Read and extract text from a Word document (.docx).

    Args:
        file_path: Path to the Word document

    Returns:
        Dict with extracted text and metadata
    """
    try:
        from docx import Document

        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        doc = Document(path)

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract tables
        tables_data = []
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                table_rows.append([cell.text for cell in row.cells])
            tables_data.append(table_rows)

        full_text = "\n\n".join(paragraphs)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables_data),
            "content": full_text,
            "tables": tables_data
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read Word document: {str(e)}"}


def read_pptx_file(file_path: str) -> dict:
    """
    Read and extract content from a PowerPoint presentation (.pptx).

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        Dict with slide content and metadata
    """
    try:
        from pptx import Presentation

        path = _resolve_file_path(file_path)

        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(path)

        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "title": "",
                "content": [],
                "notes": ""
            }

            # Get title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # Get all text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        slide_data["content"].append(shape.text)

            # Get notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data["notes"] = notes_frame.text

            slides_content.append(slide_data)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "slide_count": len(slides_content),
            "slides": slides_content
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to read PowerPoint: {str(e)}"}


def list_files(directory: str = ".", pattern: str = "*") -> dict:
    """
    List files in a directory with optional pattern matching.

    Args:
        directory: Directory path (relative to project root or absolute)
        pattern: Glob pattern (e.g., "*.json", "*.md")

    Returns:
        Dict with file listing
    """
    try:
        path = _resolve_file_path(directory)

        if not path.exists():
            return {"error": f"Directory not found: {directory}"}

        if not path.is_dir():
            return {"error": f"Not a directory: {directory}"}

        files = []
        for f in sorted(path.glob(pattern)):
            if f.is_file():
                files.append({
                    "name": f.name,
                    "path": str(f.relative_to(PROJECT_ROOT)),
                    "size_bytes": f.stat().st_size,
                    "type": f.suffix
                })

        # Also list subdirectories
        subdirs = [d.name for d in sorted(path.iterdir()) if d.is_dir() and not d.name.startswith('.')]

        return {
            "directory": str(path.relative_to(PROJECT_ROOT)),
            "pattern": pattern,
            "file_count": len(files),
            "files": files,
            "subdirectories": subdirs
        }
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"Failed to list directory: {str(e)}"}


# =============================================================================
# ACADEMY CONTENT SYNC TOOLS
# =============================================================================

def _get_file_hash(file_path: Path) -> str:
    """Get MD5 hash of a file for comparison."""
    import hashlib
    return hashlib.md5(file_path.read_bytes()).hexdigest()


def _get_file_mtime(file_path: Path) -> float:
    """Get modification time of a file."""
    return file_path.stat().st_mtime if file_path.exists() else 0


def check_academy_sync_status() -> dict:
    """
    Check sync status between local academy files and quiz-demo source of truth.

    Returns:
        Dict with sync status for curriculum.json and all .md files
    """
    try:
        if not QUIZ_DEMO_CONTENT.exists():
            return {"error": f"Quiz-demo content directory not found: {QUIZ_DEMO_CONTENT}"}

        # Ensure local directory exists
        LOCAL_ACADEMY_PATH.mkdir(parents=True, exist_ok=True)

        status = {
            "source_path": str(QUIZ_DEMO_CONTENT),
            "local_path": str(LOCAL_ACADEMY_PATH),
            "curriculum": {},
            "lessons": [],
            "summary": {
                "synced": 0,
                "local_newer": 0,
                "source_newer": 0,
                "local_only": 0,
                "source_only": 0
            }
        }

        # Check curriculum.json
        source_curriculum = QUIZ_DEMO_CONTENT / "curriculum.json"
        local_curriculum = LOCAL_ACADEMY_PATH / "curriculum.json"

        if source_curriculum.exists():
            source_mtime = _get_file_mtime(source_curriculum)
            local_mtime = _get_file_mtime(local_curriculum)

            if not local_curriculum.exists():
                sync_status = "source_only"
            elif abs(source_mtime - local_mtime) < 1:  # Within 1 second
                # Check content hash for exact match
                if _get_file_hash(source_curriculum) == _get_file_hash(local_curriculum):
                    sync_status = "synced"
                elif source_mtime > local_mtime:
                    sync_status = "source_newer"
                else:
                    sync_status = "local_newer"
            elif source_mtime > local_mtime:
                sync_status = "source_newer"
            else:
                sync_status = "local_newer"

            status["curriculum"] = {
                "source_modified": datetime.fromtimestamp(source_mtime).isoformat() if source_mtime else None,
                "local_modified": datetime.fromtimestamp(local_mtime).isoformat() if local_mtime else None,
                "status": sync_status
            }
            status["summary"][sync_status] += 1

        # Check all .md files
        source_md_files = set(f.name for f in QUIZ_DEMO_CONTENT.glob("*.md"))
        local_md_files = set(f.name for f in LOCAL_ACADEMY_PATH.glob("*.md"))

        all_md_files = source_md_files | local_md_files

        for md_file in sorted(all_md_files):
            source_file = QUIZ_DEMO_CONTENT / md_file
            local_file = LOCAL_ACADEMY_PATH / md_file

            source_exists = source_file.exists()
            local_exists = local_file.exists()

            if source_exists and not local_exists:
                sync_status = "source_only"
            elif local_exists and not source_exists:
                sync_status = "local_only"
            else:
                source_mtime = _get_file_mtime(source_file)
                local_mtime = _get_file_mtime(local_file)

                if _get_file_hash(source_file) == _get_file_hash(local_file):
                    sync_status = "synced"
                elif source_mtime > local_mtime:
                    sync_status = "source_newer"
                else:
                    sync_status = "local_newer"

            status["lessons"].append({
                "file": md_file,
                "status": sync_status
            })
            status["summary"][sync_status] += 1

        return status

    except Exception as e:
        return {"error": f"Failed to check sync status: {str(e)}"}


def sync_from_quiz_demo(files: list[str] = None, force: bool = False) -> dict:
    """
    Pull files from quiz-demo (source of truth) to local shaman-assistant.

    Args:
        files: Specific files to sync (e.g., ["curriculum.json", "clm-intro.md"]).
               If None, syncs all files that are newer in source.
        force: If True, overwrite local files even if they're newer.

    Returns:
        Dict with sync results
    """
    try:
        import shutil

        if not QUIZ_DEMO_CONTENT.exists():
            return {"error": f"Quiz-demo content directory not found: {QUIZ_DEMO_CONTENT}"}

        LOCAL_ACADEMY_PATH.mkdir(parents=True, exist_ok=True)

        results = {
            "synced": [],
            "skipped": [],
            "errors": []
        }

        # Determine which files to sync
        if files:
            files_to_check = files
        else:
            # Get all content files from source
            files_to_check = ["curriculum.json"] + [f.name for f in QUIZ_DEMO_CONTENT.glob("*.md")]

        for filename in files_to_check:
            source_file = QUIZ_DEMO_CONTENT / filename
            local_file = LOCAL_ACADEMY_PATH / filename

            if not source_file.exists():
                results["errors"].append(f"{filename}: not found in source")
                continue

            # Check if we should sync
            should_sync = force

            if not should_sync:
                if not local_file.exists():
                    should_sync = True
                elif _get_file_hash(source_file) != _get_file_hash(local_file):
                    source_mtime = _get_file_mtime(source_file)
                    local_mtime = _get_file_mtime(local_file)

                    if source_mtime > local_mtime:
                        should_sync = True
                    else:
                        results["skipped"].append(f"{filename}: local is newer (use force=True to overwrite)")
                        continue
                else:
                    results["skipped"].append(f"{filename}: already synced")
                    continue

            if should_sync:
                shutil.copy2(source_file, local_file)
                results["synced"].append(filename)

        results["summary"] = f"Synced {len(results['synced'])} files, skipped {len(results['skipped'])}, errors {len(results['errors'])}"

        return results

    except Exception as e:
        return {"error": f"Failed to sync from quiz-demo: {str(e)}"}


def sync_to_quiz_demo(files: list[str] = None, force: bool = False) -> dict:
    """
    Push local files to quiz-demo (source of truth).

    Args:
        files: Specific files to sync (e.g., ["curriculum.json", "clm-intro.md"]).
               If None, syncs all files that are newer locally.
        force: If True, overwrite source files even if they're newer.

    Returns:
        Dict with sync results
    """
    try:
        import shutil

        if not QUIZ_DEMO_CONTENT.exists():
            return {"error": f"Quiz-demo content directory not found: {QUIZ_DEMO_CONTENT}"}

        if not LOCAL_ACADEMY_PATH.exists():
            return {"error": f"Local academy directory not found: {LOCAL_ACADEMY_PATH}"}

        results = {
            "synced": [],
            "skipped": [],
            "errors": []
        }

        # Determine which files to sync
        if files:
            files_to_check = files
        else:
            # Get all content files from local
            files_to_check = []
            if (LOCAL_ACADEMY_PATH / "curriculum.json").exists():
                files_to_check.append("curriculum.json")
            files_to_check.extend([f.name for f in LOCAL_ACADEMY_PATH.glob("*.md")])

        for filename in files_to_check:
            local_file = LOCAL_ACADEMY_PATH / filename
            source_file = QUIZ_DEMO_CONTENT / filename

            if not local_file.exists():
                results["errors"].append(f"{filename}: not found locally")
                continue

            # Check if we should sync
            should_sync = force

            if not should_sync:
                if not source_file.exists():
                    should_sync = True
                elif _get_file_hash(source_file) != _get_file_hash(local_file):
                    source_mtime = _get_file_mtime(source_file)
                    local_mtime = _get_file_mtime(local_file)

                    if local_mtime > source_mtime:
                        should_sync = True
                    else:
                        results["skipped"].append(f"{filename}: source is newer (use force=True to overwrite)")
                        continue
                else:
                    results["skipped"].append(f"{filename}: already synced")
                    continue

            if should_sync:
                shutil.copy2(local_file, source_file)
                results["synced"].append(filename)

        results["summary"] = f"Synced {len(results['synced'])} files to quiz-demo, skipped {len(results['skipped'])}, errors {len(results['errors'])}"

        return results

    except Exception as e:
        return {"error": f"Failed to sync to quiz-demo: {str(e)}"}


# Tool definitions for Claude API
DOCUMENT_TOOLS = [
    {
        "name": "create_presentation",
        "description": "Create a PowerPoint presentation (.pptx). Use for customer onboarding decks, training materials, sales presentations.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "filename": {"type": "string", "description": "Output filename without extension"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "content": {"type": "array", "items": {"type": "string"}},
                            "notes": {"type": "string"},
                            "layout": {"type": "string", "enum": ["title", "section", "content", "two_column"]}
                        },
                        "required": ["title"]
                    }
                }
            },
            "required": ["title", "filename", "slides"]
        }
    },
    {
        "name": "create_document",
        "description": "Create a Word document (.docx). Use for SOPs, implementation guides, proposals, onboarding documentation.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "filename": {"type": "string"},
                "sections": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "level": {"type": "integer", "default": 1},
                            "content": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}}
                        },
                        "required": ["heading"]
                    }
                }
            },
            "required": ["title", "filename", "sections"]
        }
    },
    {
        "name": "create_pdf_report",
        "description": "Create a PDF report. Use for compliance reports, branded deliverables, formal documentation.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "filename": {"type": "string"},
                "content": {"type": "string", "description": "Markdown content"}
            },
            "required": ["title", "filename", "content"]
        }
    },
    {
        "name": "create_markdown",
        "description": "Create a Markdown file (.md). Use for documentation, README files, notes, specifications, or any text-based content with formatting.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Document title (becomes H1 heading)"},
                "filename": {"type": "string", "description": "Output filename without extension"},
                "content": {"type": "string", "description": "Markdown-formatted content (supports headings, lists, code blocks, tables, etc.)"}
            },
            "required": ["title", "filename", "content"]
        }
    },
    {
        "name": "create_json_structure",
        "description": "Create a JSON file with structured data. Use when the user provides an example structure or schema and wants to generate data in that format. Great for configuration files, data exports, API payloads, or any structured data.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "Output filename without extension"},
                "data": {
                    "description": "The JSON data structure to create (can be object or array)",
                    "oneOf": [
                        {"type": "object"},
                        {"type": "array"}
                    ]
                },
                "schema_description": {"type": "string", "description": "Optional description of the schema/structure used"}
            },
            "required": ["filename", "data"]
        }
    }
]


# File reading tool definitions for Claude API
FILE_TOOLS = [
    {
        "name": "list_files",
        "description": "List files in a directory within the project. Use to explore available files, find documents, or see what's in a folder. Supports glob patterns.",
        "input_schema": {
            "type": "object",
            "properties": {
                "directory": {
                    "type": "string",
                    "description": "Directory path relative to project root (e.g., 'output', 'output/shaman-academy', '.'). Default: project root"
                },
                "pattern": {
                    "type": "string",
                    "description": "Glob pattern to filter files (e.g., '*.json', '*.md', '*.pdf'). Default: '*' (all files)"
                }
            }
        }
    },
    {
        "name": "read_text_file",
        "description": "Read a plain text file (.txt, .md, .log, .py, .html, etc.). Returns the full content of the file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the file (relative to project root or absolute)"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_json_file",
        "description": "Read and parse a JSON file. Returns the parsed data structure.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the JSON file (relative to project root or absolute)"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_csv_file",
        "description": "Read a CSV file and return as structured data with rows and columns.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the CSV file"
                },
                "max_rows": {
                    "type": "integer",
                    "description": "Maximum rows to return (default: 100)"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_excel_file",
        "description": "Read an Excel file (.xlsx). Returns data from a specific sheet or the first sheet by default.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the Excel file"
                },
                "sheet_name": {
                    "type": "string",
                    "description": "Name of the sheet to read (default: first sheet)"
                },
                "max_rows": {
                    "type": "integer",
                    "description": "Maximum rows to return (default: 100)"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_pdf_file",
        "description": "Read and extract text from a PDF file. Extracts text content from all pages.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the PDF file"
                },
                "max_pages": {
                    "type": "integer",
                    "description": "Maximum pages to read (default: 20)"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_word_document",
        "description": "Read and extract text from a Word document (.docx). Returns paragraphs and tables.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the Word document"
                }
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_powerpoint",
        "description": "Read and extract content from a PowerPoint presentation (.pptx). Returns slide titles, content, and notes.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the PowerPoint file"
                }
            },
            "required": ["file_path"]
        }
    }
]


# Sync tool definitions for Claude API
SYNC_TOOLS = [
    {
        "name": "check_academy_sync",
        "description": "Check sync status between local academy files and the quiz-demo source of truth. Shows which files are synced, which are newer locally, and which are newer in source.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "sync_from_source",
        "description": "Pull academy content (curriculum.json and lesson .md files) from quiz-demo (source of truth) to local shaman-assistant. Use this to get the latest content.",
        "input_schema": {
            "type": "object",
            "properties": {
                "files": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Specific files to sync (e.g., ['curriculum.json', 'clm-intro.md']). If not provided, syncs all files that are newer in source."
                },
                "force": {
                    "type": "boolean",
                    "description": "If true, overwrite local files even if they're newer. Default: false"
                }
            }
        }
    },
    {
        "name": "sync_to_source",
        "description": "Push local academy content to quiz-demo (source of truth). Use this after editing files locally to update the main repository.",
        "input_schema": {
            "type": "object",
            "properties": {
                "files": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Specific files to sync (e.g., ['curriculum.json', 'clm-intro.md']). If not provided, syncs all files that are newer locally."
                },
                "force": {
                    "type": "boolean",
                    "description": "If true, overwrite source files even if they're newer. Default: false"
                }
            }
        }
    }
]


def process_document_tool(tool_name: str, tool_input: dict) -> str:
    """
    Process a document generation tool call.

    Args:
        tool_name: Name of the tool to execute
        tool_input: Tool input parameters

    Returns:
        Path to created file or error message
    """
    try:
        if tool_name == "create_presentation":
            return create_pptx(**tool_input)
        elif tool_name == "create_document":
            return create_docx(**tool_input)
        elif tool_name == "create_pdf_report":
            return create_pdf(**tool_input)
        elif tool_name == "create_markdown":
            return create_markdown(**tool_input)
        elif tool_name == "create_json_structure":
            return create_json(**tool_input)
        else:
            return f"Unknown document tool: {tool_name}"
    except Exception as e:
        return f"Error creating document: {str(e)}"


def process_file_tool(tool_name: str, tool_input: dict) -> dict:
    """
    Process a file reading tool call.

    Args:
        tool_name: Name of the tool to execute
        tool_input: Tool input parameters

    Returns:
        Dict with file content/data or error
    """
    try:
        if tool_name == "list_files":
            return list_files(
                directory=tool_input.get("directory", "."),
                pattern=tool_input.get("pattern", "*")
            )
        elif tool_name == "read_text_file":
            return read_text_file(tool_input["file_path"])
        elif tool_name == "read_json_file":
            return read_json_file(tool_input["file_path"])
        elif tool_name == "read_csv_file":
            return read_csv_file(
                tool_input["file_path"],
                max_rows=tool_input.get("max_rows", 100)
            )
        elif tool_name == "read_excel_file":
            return read_excel_file(
                tool_input["file_path"],
                sheet_name=tool_input.get("sheet_name"),
                max_rows=tool_input.get("max_rows", 100)
            )
        elif tool_name == "read_pdf_file":
            return read_pdf_file(
                tool_input["file_path"],
                max_pages=tool_input.get("max_pages", 20)
            )
        elif tool_name == "read_word_document":
            return read_docx_file(tool_input["file_path"])
        elif tool_name == "read_powerpoint":
            return read_pptx_file(tool_input["file_path"])
        else:
            return {"error": f"Unknown file tool: {tool_name}"}
    except Exception as e:
        return {"error": f"Error reading file: {str(e)}"}


def process_sync_tool(tool_name: str, tool_input: dict) -> dict:
    """
    Process an academy sync tool call.

    Args:
        tool_name: Name of the tool to execute
        tool_input: Tool input parameters

    Returns:
        Dict with sync results or error
    """
    try:
        if tool_name == "check_academy_sync":
            return check_academy_sync_status()
        elif tool_name == "sync_from_source":
            return sync_from_quiz_demo(
                files=tool_input.get("files"),
                force=tool_input.get("force", False)
            )
        elif tool_name == "sync_to_source":
            return sync_to_quiz_demo(
                files=tool_input.get("files"),
                force=tool_input.get("force", False)
            )
        else:
            return {"error": f"Unknown sync tool: {tool_name}"}
    except Exception as e:
        return {"error": f"Error in sync operation: {str(e)}"}
